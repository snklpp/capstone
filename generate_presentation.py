#!/usr/bin/env python3
"""Generate a detailed project presentation PPTX for the IIT Hyderabad SSI Portal.

Follows the structure and style of Introduction-to-System-Requirement-Organizations_Bilingual.pptx:
  - 10×5.62 inch slides
  - SRS sections 1-10 + architecture diagrams
  - Clean formatting with diagrams, tables, and callout boxes
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ── Slide dimensions (matching reference) ────────────────────────────────────
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

# ── Color palette ─────────────────────────────────────────────────────────────
DARK_BLUE = RGBColor(0x0D, 0x47, 0xA1)
MED_BLUE = RGBColor(0x1A, 0x23, 0x7E)
LIGHT_BLUE = RGBColor(0x42, 0xA5, 0xF5)
ACCENT_GREEN = RGBColor(0x1B, 0x5E, 0x20)
ACCENT_PURPLE = RGBColor(0x4A, 0x14, 0x8C)
ACCENT_ORANGE = RGBColor(0xE6, 0x51, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
TABLE_HEADER_BG = RGBColor(0x0D, 0x47, 0xA1)
TABLE_ALT_BG = RGBColor(0xE3, 0xF2, 0xFD)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helper functions ──────────────────────────────────────────────────────────

def add_bg(slide, color=MED_BLUE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=16, bold=False,
                color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, font_size=14,
                   color=DARK_GRAY, font_name="Calibri", line_spacing=1.2, bullet=False):
    """Add a textbox with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = ("• " + line) if bullet else line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(4)
    return txBox

def add_callout_box(slide, left, top, width, height, title, body_lines,
                    fill_color=RGBColor(0xE8, 0xEA, 0xF6), border_color=MED_BLUE,
                    title_color=MED_BLUE, body_color=DARK_GRAY):
    """Rounded rectangle callout (like reference)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.font.name = "Calibri"

    for line in body_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = body_color
        p.font.name = "Calibri"
        p.space_after = Pt(2)
    return shape

def add_table(slide, left, top, width, rows_data, col_widths_pct=None):
    """Add a styled table. rows_data is list of lists."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0]) if rows_data else 1
    row_h = Inches(0.35)
    table_h = row_h * n_rows
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, table_h)
    table = shape.table

    if col_widths_pct:
        for ci, pct in enumerate(col_widths_pct):
            table.columns[ci].width = int(width * pct)

    for ri, row_data in enumerate(rows_data):
        for ci, cell_text in enumerate(row_data):
            cell = table.cell(ri, ci)
            cell.text = str(cell_text)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = "Calibri"
                if ri == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                else:
                    p.font.color.rgb = DARK_GRAY
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
    return shape

def title_slide(title_text, subtitle_text, bg_color=MED_BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_bg(slide, bg_color)
    add_textbox(slide, Inches(0.8), Inches(1.4), Inches(8.4), Inches(1.2),
                title_text, font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.8), Inches(2.8), Inches(8.4), Inches(1.0),
                subtitle_text, font_size=16, color=RGBColor(0xBB, 0xDE, 0xFB), alignment=PP_ALIGN.CENTER)
    return slide

def section_header(title_text, subtitle_text=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BLUE)
    add_textbox(slide, Inches(0.8), Inches(1.8), Inches(8.4), Inches(1.0),
                title_text, font_size=30, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    if subtitle_text:
        add_textbox(slide, Inches(0.8), Inches(3.0), Inches(8.4), Inches(0.8),
                    subtitle_text, font_size=14, color=RGBColor(0x90, 0xCA, 0xF9), alignment=PP_ALIGN.CENTER)
    return slide

def content_slide(title_text, bg_color=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if bg_color != WHITE:
        add_bg(slide, bg_color)
    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()
    add_textbox(slide, Inches(0.5), Inches(0.1), Inches(9), Inches(0.6),
                title_text, font_size=24, bold=True, color=WHITE)
    return slide


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ══════════════════════════════════════════════════════════════════════════════

s = title_slide(
    "IIT Hyderabad — SSI University Credential Portal",
    "Self-Sovereign Identity with Decentralized Identifiers (DIDs)\n"
    "and Verifiable Credentials (VCs)\n\n"
    "Software Requirements Specification & System Design"
)
add_textbox(s, Inches(0.8), Inches(4.2), Inches(8.4), Inches(0.6),
            "Capstone Project  •  FastAPI + W3C Standards + OID4VCI/VP  •  2025",
            font_size=12, color=RGBColor(0x90, 0xCA, 0xF9), alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Table of Contents
# ══════════════════════════════════════════════════════════════════════════════

s = content_slide("Table of Contents")
toc = [
    "1.  System Overview & Problem Statement",
    "2.  Stakeholders & User Roles",
    "3.  System Architecture (Block Diagram)",
    "4.  Functional Requirements",
    "5.  Non-Functional Requirements",
    "6.  Data Requirements & ER Model",
    "7.  External Interfaces",
    "8.  Credential Issuance Flow (OID4VCI)",
    "9.  Verification Flow (OID4VP)",
    "10. Constraints & Assumptions",
    "11. Acceptance Criteria & Traceability",
    "12. Technology Stack",
    "13. UI Design (Screenshots)",
    "14. Reference Standards",
]
add_multi_text(s, Inches(1.0), Inches(1.1), Inches(8.0), Inches(4.0),
               toc, font_size=15, color=DARK_GRAY, line_spacing=1.4)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Section Header: System Overview
# ══════════════════════════════════════════════════════════════════════════════

section_header("1. System Overview", "Problem Statement & Objectives")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — System Overview Content
# ══════════════════════════════════════════════════════════════════════════════

s = content_slide("System Overview")

add_callout_box(s, Inches(0.4), Inches(1.0), Inches(4.4), Inches(2.0),
    "Problem Statement",
    [
        "Traditional university credentials are paper-based,",
        "easy to forge, and hard to verify instantly.",
        "Employers and third parties have no standard",
        "machine-readable way to confirm authenticity.",
    ],
    fill_color=RGBColor(0xFF, 0xEB, 0xEE), border_color=RGBColor(0xC6, 0x28, 0x28),
    title_color=RGBColor(0xC6, 0x28, 0x28))

add_callout_box(s, Inches(5.2), Inches(1.0), Inches(4.4), Inches(2.0),
    "Solution",
    [
        "A Self-Sovereign Identity (SSI) portal where",
        "IIT Hyderabad issues W3C Verifiable Credentials",
        "to student wallets (Altme). Verifiers can request",
        "cryptographic proofs without contacting the issuer.",
    ],
    fill_color=RGBColor(0xE8, 0xF5, 0xE9), border_color=ACCENT_GREEN,
    title_color=ACCENT_GREEN)

add_callout_box(s, Inches(0.4), Inches(3.3), Inches(9.2), Inches(1.9),
    "Key Objectives",
    [
        "✓  Issue tamper-proof digital credentials (Degree, Internship, Skill Badge) via OID4VCI",
        "✓  Enable students to store credentials in a mobile wallet (Altme) — full data sovereignty",
        "✓  Allow any verifier to request & verify presentations (OID4VP) — no issuer in the loop",
        "✓  Implement W3C DID:web identifiers with key rotation and key history",
    ],
    fill_color=RGBColor(0xE8, 0xEA, 0xF6), border_color=MED_BLUE,
    title_color=MED_BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Section Header: Stakeholders
# ══════════════════════════════════════════════════════════════════════════════

section_header("2. Stakeholders & User Roles", "Actors, Responsibilities, and Interactions")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Stakeholders Table
# ══════════════════════════════════════════════════════════════════════════════

s = content_slide("Stakeholders & User Roles")

stakeholder_data = [
    ["Role", "Actor", "Responsibilities", "Access Level"],
    ["Student", "IIT Hyderabad Students",
     "Register, create DID, receive VCs, present VPs",
     "Own credentials"],
    ["Admin (Issuer)", "University Administration",
     "View students, issue credentials (Degree/Internship/Skill)",
     "Credential issuance"],
    ["Verifier", "Employers / Third Parties",
     "Initiate OID4VP requests, verify presentations",
     "Read-only verification"],
    ["Wallet", "Altme Mobile Wallet",
     "Scan OID4VCI offers, store VCs, respond to OID4VP",
     "Student's device"],
]
add_table(s, Inches(0.4), Inches(1.1), Inches(9.2), stakeholder_data,
          col_widths_pct=[0.12, 0.2, 0.45, 0.23])


add_callout_box(s, Inches(0.4), Inches(3.3), Inches(9.2), Inches(1.8),
    "Trust Triangle (SSI Model)",
    [
        "Issuer (Admin)  ——issues VC——▶  Holder (Student Wallet)  ◀——verifies VP——  Verifier",
        "• Issuer signs credentials with its ES256 private key → Holder stores JWT in wallet",
        "• Verifier never contacts Issuer directly — resolves DID Document to get public key",
        "• Holder controls what to disclose — presents Verifiable Presentation containing VC(s)",
    ],
    fill_color=RGBColor(0xFD, 0xF1, 0xE0), border_color=ACCENT_ORANGE,
    title_color=ACCENT_ORANGE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Section Header: Architecture
# ══════════════════════════════════════════════════════════════════════════════

section_header("3. System Architecture", "Block Diagram & Component Overview")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Block Diagram (drawn with shapes)
# ══════════════════════════════════════════════════════════════════════════════

s = content_slide("System Architecture — Block Diagram")

def draw_block(slide, left, top, w, h, text, fill_c, text_c=WHITE, font_sz=10):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_c
    shape.line.color.rgb = RGBColor(0x90, 0x90, 0x90)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_sz)
    p.font.bold = True
    p.font.color.rgb = text_c
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape

def draw_arrow_label(slide, left, top, w, h, text, font_sz=8):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_sz)
    p.font.color.rgb = DARK_GRAY
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

# Top row: Frontend
draw_block(s, Inches(3.3), Inches(1.0), Inches(3.4), Inches(0.55),
           "SPA Frontend (HTML / JS / CSS)", DARK_BLUE, WHITE, 11)

# Middle row: FastAPI + routers
draw_block(s, Inches(0.3), Inches(2.0), Inches(1.7), Inches(0.55),
           "Auth Router\n/auth/*", RGBColor(0x38, 0x8E, 0x3C), WHITE, 9)
draw_block(s, Inches(2.2), Inches(2.0), Inches(1.7), Inches(0.55),
           "Student Router\n/students/*", RGBColor(0x15, 0x65, 0xC0), WHITE, 9)
draw_block(s, Inches(4.1), Inches(2.0), Inches(1.7), Inches(0.55),
           "Admin Router\n/admin/*", RGBColor(0x6A, 0x1B, 0x9A), WHITE, 9)
draw_block(s, Inches(6.0), Inches(2.0), Inches(1.7), Inches(0.55),
           "Verify Router\n/verify/*", ACCENT_ORANGE, WHITE, 9)
draw_block(s, Inches(7.9), Inches(2.0), Inches(1.7), Inches(0.55),
           "DID Router\n/.well-known/*", RGBColor(0x00, 0x69, 0x5C), WHITE, 9)

# Label
draw_arrow_label(s, Inches(3.3), Inches(1.6), Inches(3.4), Inches(0.35),
                 "▼  REST API  (FastAPI)  ▼", 9)

# Bottom row: Services
draw_block(s, Inches(0.3), Inches(3.1), Inches(2.0), Inches(0.55),
           "Crypto Module\nES256 / JWK / JWT", RGBColor(0x4E, 0x34, 0x2E), WHITE, 9)
draw_block(s, Inches(2.5), Inches(3.1), Inches(2.0), Inches(0.55),
           "DID Resolver\ndid:web / did:jwk / did:key", RGBColor(0x00, 0x57, 0x8A), WHITE, 9)
draw_block(s, Inches(4.7), Inches(3.1), Inches(2.0), Inches(0.55),
           "SQLite Database\nORM: SQLAlchemy", RGBColor(0x54, 0x4F, 0x4D), WHITE, 9)
draw_block(s, Inches(6.9), Inches(3.1), Inches(2.7), Inches(0.55),
           "Cloudflare Tunnel\n(Public HTTPS Endpoint)", RGBColor(0xF5, 0x7C, 0x00), WHITE, 9)

# External actors
draw_block(s, Inches(0.3), Inches(4.2), Inches(2.3), Inches(0.55),
           "🎓  Student Browser", RGBColor(0xE3, 0xF2, 0xFD), DARK_BLUE, 10)
draw_block(s, Inches(2.9), Inches(4.2), Inches(2.0), Inches(0.55),
           "📱  Altme Wallet", RGBColor(0xE8, 0xF5, 0xE9), ACCENT_GREEN, 10)
draw_block(s, Inches(5.2), Inches(4.2), Inches(2.3), Inches(0.55),
           "🏢  Verifier Dashboard", RGBColor(0xFD, 0xF1, 0xE0), ACCENT_ORANGE, 10)
draw_block(s, Inches(7.8), Inches(4.2), Inches(1.8), Inches(0.55),
           "🔑  Issuer (Admin)", RGBColor(0xF3, 0xE5, 0xF5), ACCENT_PURPLE, 10)

draw_arrow_label(s, Inches(3.0), Inches(2.65), Inches(4.0), Inches(0.35),
                 "▼  Internal Services  ▼", 9)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Section Header: Functional Requirements
# ══════════════════════════════════════════════════════════════════════════════

section_header("4. Functional Requirements", "Core features by module")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Functional Requirements: Auth & Student
# ══════════════════════════════════════════════════════════════════════════════

s = content_slide("Functional Requirements — Authentication & Student Module")

fr_auth = [
    ["ID", "Requirement", "Endpoint", "Method"],
    ["FR-01", "Student registration with roll number, name, password", "/auth/register", "POST"],
    ["FR-02", "Login with JWT token issuance (HS256, 60 min expiry)", "/auth/login", "POST"],
    ["FR-03", "OID4VCI token exchange (pre-authorized code → access token)", "/auth/token", "POST"],
    ["FR-04", "Role-based access control (student / admin / verifier)", "Middleware", "—"],
]
add_table(s, Inches(0.3), Inches(1.0), Inches(9.4), fr_auth,
          col_widths_pct=[0.08, 0.47, 0.3, 0.15])

fr_student = [
    ["ID", "Requirement", "Endpoint", "Method"],
    ["FR-05", "Create DID (did:web) with EC P-256 public key", "/students/did", "POST"],
    ["FR-06", "Key rotation with archival of previous keys", "/students/did/rotate-key", "POST"],
    ["FR-07", "List all VCs + pending credential offers with QR codes", "/students/vcs", "GET"],
    ["FR-08", "View pending verification challenges targeting student", "/students/challenges", "GET"],
]
add_table(s, Inches(0.3), Inches(3.0), Inches(9.4), fr_student,
          col_widths_pct=[0.08, 0.47, 0.3, 0.15])


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Functional Requirements: Admin & Verify
# ══════════════════════════════════════════════════════════════════════════════

s = content_slide("Functional Requirements — Admin & Verification Module")

fr_admin = [
    ["ID", "Requirement", "Endpoint", "Method"],
    ["FR-09", "List all students with DID status", "/admin/students", "GET"],
    ["FR-10", "Create a credential offer (pre-auth code + QR)", "/admin/issue-vc", "POST"],
    ["FR-11", "OID4VCI credential endpoint — signs VC on wallet claim", "/admin/issue", "POST"],
    ["FR-12", "Serve OID4VCI issuer metadata", "/.well-known/openid-credential-issuer", "GET"],
]
add_table(s, Inches(0.3), Inches(1.0), Inches(9.4), fr_admin,
          col_widths_pct=[0.08, 0.47, 0.3, 0.15])

fr_verify = [
    ["ID", "Requirement", "Endpoint", "Method"],
    ["FR-13", "Initiate OID4VP verification session with QR", "/verify/init", "POST"],
    ["FR-14", "Serve signed OID4VP request object (JWT)", "/api/verify/request/{id}", "GET"],
    ["FR-15", "Receive wallet VP via direct_post callback", "/api/verify/callback/{id}", "POST"],
    ["FR-16", "Poll verification session status", "/verify/status/{id}", "GET"],
    ["FR-17", "List verification history for verifier", "/verify/history", "GET"],
]
add_table(s, Inches(0.3), Inches(3.0), Inches(9.4), fr_verify,
          col_widths_pct=[0.08, 0.47, 0.3, 0.15])


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Section Header: Non-Functional Requirements
# ══════════════════════════════════════════════════════════════════════════════

section_header("5. Non-Functional Requirements", "Security, Performance, Usability, Maintainability")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Non-Functional Requirements Content
# ══════════════════════════════════════════════════════════════════════════════

s = content_slide("Non-Functional Requirements")

add_callout_box(s, Inches(0.3), Inches(1.0), Inches(4.5), Inches(1.8),
    "🔒  Security",
    [
        "• Bcrypt password hashing (passlib)",
        "• JWT bearer tokens (HS256) with expiry",
        "• ECDSA P-256 (ES256) for VC/VP signing",
        "• Role-based endpoint guards",
        "• CORS middleware for cross-origin control",
    ],
    fill_color=RGBColor(0xFC, 0xE4, 0xEC), border_color=RGBColor(0xAD, 0x14, 0x57),
    title_color=RGBColor(0xAD, 0x14, 0x57))

add_callout_box(s, Inches(5.2), Inches(1.0), Inches(4.5), Inches(1.8),
    "⚡  Performance & Scalability",
    [
        "• SQLite for dev; pluggable to PostgreSQL",
        "• Stateless JWT auth (no server sessions)",
        "• Lightweight FastAPI (async capable)",
        "• Pre-authorized code flow avoids real-time signing",
        "• Credential offers stored for async wallet claim",
    ],
    fill_color=RGBColor(0xE0, 0xF7, 0xFA), border_color=RGBColor(0x00, 0x69, 0x5C),
    title_color=RGBColor(0x00, 0x69, 0x5C))

add_callout_box(s, Inches(0.3), Inches(3.1), Inches(4.5), Inches(1.8),
    "🎨  Usability",
    [
        "• Single-page app with focus-mode navigation",
        "• Hamburger menu with role-specific dashboards",
        "• QR codes for wallet scanning (OID4VCI + OID4VP)",
        "• Responsive design for desktop & mobile",
        "• Real-time polling for verification status",
    ],
    fill_color=RGBColor(0xE8, 0xEA, 0xF6), border_color=MED_BLUE,
    title_color=MED_BLUE)

add_callout_box(s, Inches(5.2), Inches(3.1), Inches(4.5), Inches(1.8),
    "🔧  Maintainability & Portability",
    [
        "• Modular router architecture (auth, student, admin, verify, DID)",
        "• Pydantic schemas for all request/response validation",
        "• SQLAlchemy ORM with declarative models",
        "• Persistent issuer key on disk (survives restarts)",
        "• Cloudflare tunnel for easy public exposure",
    ],
    fill_color=RGBColor(0xF3, 0xE5, 0xF5), border_color=ACCENT_PURPLE,
    title_color=ACCENT_PURPLE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Section Header: Data Requirements
# ══════════════════════════════════════════════════════════════════════════════

section_header("6. Data Requirements", "Entity-Relationship Model & Database Schema")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — ER Diagram (drawn with shapes)
# ══════════════════════════════════════════════════════════════════════════════

s = content_slide("Data Model — Entity Relationship Diagram")

def er_entity(slide, left, top, w, h, name, attrs, fill_c):
    # Entity header
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, Inches(0.38))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = fill_c
    hdr.line.color.rgb = fill_c
    tf = hdr.text_frame
    tf.margin_left = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    # Attributes box
    body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + Inches(0.38), w, h - Inches(0.38))
    body.fill.solid()
    body.fill.fore_color.rgb = WHITE
    body.line.color.rgb = fill_c
    body.line.width = Pt(1.5)
    tf = body.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_top = Inches(0.04)
    for i, attr in enumerate(attrs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = attr
        p.font.size = Pt(8)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Calibri"
        p.space_after = Pt(1)

# Users
er_entity(s, Inches(0.3), Inches(1.0), Inches(1.8), Inches(1.6), "Users", [
    "PK  id (UUID)", "username (unique)", "hashed_password",
    "role (enum)", "full_name", "student_id", "created_at"
], DARK_BLUE)

# DIDDocument
er_entity(s, Inches(2.6), Inches(1.0), Inches(1.8), Inches(1.6), "DIDDocument", [
    "PK  id (UUID)", "did_uri (unique)", "FK  user_id",
    "public_key_jwk (JSON)", "did_document (JSON)", "key_version", "created_at"
], RGBColor(0x15, 0x65, 0xC0))

# VerifiableCredential
er_entity(s, Inches(4.9), Inches(1.0), Inches(1.8), Inches(1.6), "VerifiableCredential", [
    "PK  id (UUID)", "vc_id (unique)", "FK  did_id",
    "vc_type", "vc_jwt (TEXT)", "issued_at"
], ACCENT_GREEN)

# CredentialOffer
er_entity(s, Inches(7.1), Inches(1.0), Inches(1.8), Inches(1.6), "CredentialOffer", [
    "PK  id (UUID)", "FK  user_id", "vc_type",
    "offer_details (JSON)", "pre_authorized_code", "status (enum)", "created_at"
], ACCENT_PURPLE)

# VerificationSession
er_entity(s, Inches(0.3), Inches(3.2), Inches(2.4), Inches(1.8), "VerificationSession", [
    "PK  id (UUID)", "verification_id", "FK  verifier_id",
    "target_did", "FK  holder_did", "nonce", "status (enum)",
    "presentation_definition (JSON)", "vc_jwt", "vp_jwt", "expires_at"
], ACCENT_ORANGE)

# KeyHistory
er_entity(s, Inches(3.2), Inches(3.2), Inches(1.8), Inches(1.3), "KeyHistory", [
    "PK  id (UUID)", "FK  did_id",
    "public_key_jwk (JSON)", "key_id", "rotated_at"
], RGBColor(0x54, 0x4F, 0x4D))

# Relationship labels
draw_arrow_label(s, Inches(2.1), Inches(1.5), Inches(0.6), Inches(0.3), "1:1", 9)
draw_arrow_label(s, Inches(4.4), Inches(1.5), Inches(0.6), Inches(0.3), "1:N", 9)
draw_arrow_label(s, Inches(6.7), Inches(1.5), Inches(0.5), Inches(0.3), "1:N", 9)
draw_arrow_label(s, Inches(5.2), Inches(3.6), Inches(1.8), Inches(0.3), "Relationships →", 8)

# Legend
add_callout_box(s, Inches(5.5), Inches(3.3), Inches(4.1), Inches(1.6),
    "Relationships",
    [
        "Users  ←1:1→  DIDDocument   (one DID per student)",
        "DIDDocument  ←1:N→  VerifiableCredential",
        "DIDDocument  ←1:N→  KeyHistory",
        "Users  ←1:N→  CredentialOffer",
        "Users  ←1:N→  VerificationSession  (as verifier)",
        "DIDDocument  ←1:N→  VerificationSession  (as holder)",
    ])


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Section Header: External Interfaces
# ══════════════════════════════════════════════════════════════════════════════

section_header("7. External Interfaces", "User Interface, Hardware, Software, Communication")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — External Interfaces Content
# ══════════════════════════════════════════════════════════════════════════════

s = content_slide("External Interfaces")

intf_data = [
    ["Interface", "Type", "Description", "Protocol / Format"],
    ["Web Browser (SPA)", "User Interface",
     "Single-page app: Login, Student, Admin, Verifier dashboards",
     "HTML/JS/CSS over HTTPS"],
    ["Altme Wallet", "Software Interface",
     "Mobile wallet scans QR for credential offer & verification",
     "OID4VCI / OID4VP (deep links)"],
    ["OID4VCI Metadata", "Communication",
     "/.well-known/openid-credential-issuer discovery endpoint",
     "JSON over HTTPS"],
    ["OID4VP Request", "Communication",
     "Signed JWT authorization request served to wallet",
     "oauth-authz-req+jwt"],
    ["DID Resolution", "Software Interface",
     "Resolve did:web, did:jwk, did:key to DID Documents",
     "Internal resolver + HTTP"],
    ["Cloudflare Tunnel", "Communication",
     "Exposes localhost:8000 as public HTTPS for wallets",
     "cloudflared binary"],
    ["SQLite Database", "Data Interface",
     "Persistent storage for users, DIDs, VCs, sessions",
     "SQLAlchemy ORM"],
]
add_table(s, Inches(0.3), Inches(1.0), Inches(9.4), intf_data,
          col_widths_pct=[0.18, 0.14, 0.43, 0.25])


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Section Header: OID4VCI Flow
# ══════════════════════════════════════════════════════════════════════════════

section_header("8. Credential Issuance Flow", "OID4VCI — Pre-Authorized Code Grant")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — OID4VCI Sequence / Timeline Diagram
# ══════════════════════════════════════════════════════════════════════════════

s = content_slide("OID4VCI Issuance — Sequence Diagram")

# Actors
actors = [
    (Inches(0.5), "Admin\n(Issuer)", ACCENT_PURPLE),
    (Inches(2.7), "FastAPI\nBackend", DARK_BLUE),
    (Inches(5.0), "Student\nDashboard", RGBColor(0x15, 0x65, 0xC0)),
    (Inches(7.3), "Altme\nWallet", ACCENT_GREEN),
]
for x, label, color in actors:
    draw_block(s, x, Inches(1.0), Inches(1.5), Inches(0.5), label, color, WHITE, 9)
    # Vertical lifeline
    line_shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     x + Inches(0.72), Inches(1.5), Inches(0.06), Inches(3.4))
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = RGBColor(0xBD, 0xBD, 0xBD)
    line_shape.line.fill.background()

# Sequence steps
steps = [
    (0, 1, "1. POST /admin/issue-vc", Inches(1.65)),
    (1, 1, "2. Create CredentialOffer + pre-auth code", Inches(1.95)),
    (1, 2, "3. Return QR (openid-credential-offer://)", Inches(2.25)),
    (2, 3, "4. Student scans QR → Wallet opens", Inches(2.55)),
    (3, 1, "5. POST /auth/token (exchange code)", Inches(2.85)),
    (1, 3, "6. Return access_token (JWT)", Inches(3.15)),
    (3, 1, "7. POST /admin/issue (with proof JWT)", Inches(3.45)),
    (1, 1, "8. Sign VC (ES256), update DID doc", Inches(3.75)),
    (1, 3, "9. Return signed VC JWT → stored in wallet", Inches(4.05)),
]

for src_idx, dst_idx, label, y_pos in steps:
    src_x = actors[src_idx][0] + Inches(0.75)
    dst_x = actors[dst_idx][0] + Inches(0.75)
    if src_idx == dst_idx:
        # Self-message
        draw_arrow_label(s, src_x + Inches(0.1), y_pos, Inches(2.2), Inches(0.25), label, 8)
    else:
        min_x = min(src_x, dst_x)
        width = abs(int(dst_x) - int(src_x))
        draw_arrow_label(s, min_x, y_pos, Emu(width), Inches(0.25), label, 8)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — OID4VCI Credential Types
# ══════════════════════════════════════════════════════════════════════════════

s = content_slide("Supported Credential Types")

cred_types = [
    ["Credential Type", "Format", "Key Claims", "Display Color"],
    ["UniversityDegreeCredential", "jwt_vc_json",
     "student_name, student_id, degree, branch, CGPA, graduation_year, honours",
     "#1A237E (Navy)"],
    ["InternshipCredential", "jwt_vc_json",
     "student_name, company, role, duration",
     "#1b5e20 (Green)"],
    ["SkillBadgeCredential", "jwt_vc_json",
     "student_name, skill_name, proficiency",
     "#4a148c (Purple)"],
]
add_table(s, Inches(0.3), Inches(1.0), Inches(9.4), cred_types,
          col_widths_pct=[0.25, 0.12, 0.45, 0.18])

add_callout_box(s, Inches(0.3), Inches(2.8), Inches(9.4), Inches(2.4),
    "OID4VCI Protocol Details",
    [
        "Grant Type:  urn:ietf:params:oauth:grant-type:pre-authorized_code",
        "Signing Algorithm:  ES256 (ECDSA P-256)",
        "Key Binding:  Wallet proof JWT → extract public key → update student DID document",
        "Discovery:  /.well-known/openid-credential-issuer (JSON metadata)",
        "Token Endpoint:  /auth/token (exchanges pre-authorized code for access_token)",
        "Credential Endpoint:  /admin/issue (signs & returns VC JWT)",
        "Wallet Compatibility:  Altme (tested), any OID4VCI-compliant wallet",
    ])


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — Section Header: OID4VP Flow
# ══════════════════════════════════════════════════════════════════════════════

section_header("9. Verification Flow", "OID4VP — Verifiable Presentation Request")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — OID4VP Sequence Diagram
# ══════════════════════════════════════════════════════════════════════════════

s = content_slide("OID4VP Verification — Sequence Diagram")

actors_v = [
    (Inches(0.5), "Verifier\nDashboard", ACCENT_ORANGE),
    (Inches(2.7), "FastAPI\nBackend", DARK_BLUE),
    (Inches(5.0), "Altme\nWallet", ACCENT_GREEN),
    (Inches(7.3), "DID\nResolver", RGBColor(0x00, 0x69, 0x5C)),
]
for x, label, color in actors_v:
    draw_block(s, x, Inches(1.0), Inches(1.5), Inches(0.5), label, color, WHITE, 9)
    line_shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     x + Inches(0.72), Inches(1.5), Inches(0.06), Inches(3.4))
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = RGBColor(0xBD, 0xBD, 0xBD)
    line_shape.line.fill.background()

steps_v = [
    (0, 1, "1. POST /verify/init", Inches(1.65)),
    (1, 1, "2. Create session + nonce + PD", Inches(1.95)),
    (1, 0, "3. Return QR (openid4vp://...)", Inches(2.25)),
    (0, 2, "4. Wallet scans QR", Inches(2.55)),
    (2, 1, "5. GET /api/verify/request/{id} (JWT)", Inches(2.85)),
    (2, 2, "6. Select matching VC, build VP", Inches(3.15)),
    (2, 1, "7. POST /api/verify/callback/{id} (vp_token)", Inches(3.45)),
    (1, 3, "8. Resolve holder DID → get public key", Inches(3.75)),
    (1, 1, "9. Verify VP sig, VC sig → VERIFIED", Inches(4.05)),
]

for src_idx, dst_idx, label, y_pos in steps_v:
    src_x = actors_v[src_idx][0] + Inches(0.75)
    dst_x = actors_v[dst_idx][0] + Inches(0.75)
    if src_idx == dst_idx:
        draw_arrow_label(s, src_x + Inches(0.1), y_pos, Inches(2.2), Inches(0.25), label, 8)
    else:
        min_x = min(src_x, dst_x)
        width = abs(int(dst_x) - int(src_x))
        draw_arrow_label(s, min_x, y_pos, Emu(width), Inches(0.25), label, 8)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 23 — OID4VP Protocol Details
# ══════════════════════════════════════════════════════════════════════════════

s = content_slide("OID4VP Protocol Details")

add_callout_box(s, Inches(0.3), Inches(1.0), Inches(4.5), Inches(2.0),
    "Presentation Definition (DIF format)",
    [
        "format: jwt_vc_json + jwt_vp_json (alg: ES256)",
        "input_descriptors: 'iith_credential'",
        "constraints.fields.path: $.vc.type[*]",
        "filter: UniversityDegreeCredential |",
        "        InternshipCredential |",
        "        SkillBadgeCredential",
    ],
    fill_color=RGBColor(0xE8, 0xEA, 0xF6), border_color=MED_BLUE,
    title_color=MED_BLUE)

add_callout_box(s, Inches(5.2), Inches(1.0), Inches(4.5), Inches(2.0),
    "Request Object (signed JWT)",
    [
        "client_id_scheme: redirect_uri",
        "response_type: vp_token",
        "response_mode: direct_post",
        "response_uri: /api/verify/callback/{id}",
        "Signed with verifier ES256 key",
        "Header: typ=oauth-authz-req+jwt",
    ],
    fill_color=RGBColor(0xFD, 0xF1, 0xE0), border_color=ACCENT_ORANGE,
    title_color=ACCENT_ORANGE)

add_callout_box(s, Inches(0.3), Inches(3.3), Inches(4.5), Inches(1.9),
    "Verification Steps (Backend)",
    [
        "1. Extract vp_token from wallet POST",
        "2. Decode VP JWT → get holder DID from 'iss'",
        "3. Resolve holder DID → obtain public key",
        "4. Verify VP signature with holder key",
        "5. Extract embedded VC → resolve issuer DID",
        "6. Verify VC signature with issuer key",
        "7. Check nonce match → mark VERIFIED",
    ],
    fill_color=RGBColor(0xE8, 0xF5, 0xE9), border_color=ACCENT_GREEN,
    title_color=ACCENT_GREEN)

add_callout_box(s, Inches(5.2), Inches(3.3), Inches(4.5), Inches(1.9),
    "Supported DID Methods for Resolution",
    [
        "did:web — HTTP-based, our student/issuer DIDs",
        "did:jwk — Self-describing, JWK in the DID itself",
        "did:key — Compressed P-256 key (multicodec)",
        "",
        "Wallet may use any method — resolver adapts",
        "accordingly via app/resolver.py",
    ],
    fill_color=RGBColor(0xF3, 0xE5, 0xF5), border_color=ACCENT_PURPLE,
    title_color=ACCENT_PURPLE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 24 — Section Header: Constraints
# ══════════════════════════════════════════════════════════════════════════════

section_header("10. Constraints & Assumptions", "Design Constraints, Dependencies, and Limitations")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 25 — Constraints & Assumptions Content
# ══════════════════════════════════════════════════════════════════════════════

s = content_slide("Constraints & Assumptions")

add_callout_box(s, Inches(0.3), Inches(1.0), Inches(4.5), Inches(2.0),
    "Design Constraints",
    [
        "• Must use W3C VC Data Model (jwt_vc_json)",
        "• Must support OID4VCI Draft 13+ protocol",
        "• Must support OID4VP with direct_post response",
        "• ES256 (P-256) mandatory for all signing",
        "• Single-server deployment (no microservices)",
        "• Python 3.10+ required (FastAPI + type hints)",
    ],
    fill_color=RGBColor(0xFC, 0xE4, 0xEC), border_color=RGBColor(0xAD, 0x14, 0x57),
    title_color=RGBColor(0xAD, 0x14, 0x57))

add_callout_box(s, Inches(5.2), Inches(1.0), Inches(4.5), Inches(2.0),
    "Assumptions",
    [
        "• Students have a smartphone with Altme wallet",
        "• Cloudflare tunnel provides stable HTTPS URL",
        "• Admin has pre-verified student identity",
        "• Wallet implements OID4VCI + OID4VP correctly",
        "• Browser supports modern JS (ES6+)",
        "• Development uses SQLite (prod: PostgreSQL)",
    ],
    fill_color=RGBColor(0xE0, 0xF7, 0xFA), border_color=RGBColor(0x00, 0x69, 0x5C),
    title_color=RGBColor(0x00, 0x69, 0x5C))

add_callout_box(s, Inches(0.3), Inches(3.3), Inches(9.4), Inches(1.9),
    "Dependencies",
    [
        "• FastAPI 0.100+ (web framework)  •  SQLAlchemy 2.x (ORM)  •  python-jose (JWT)",
        "• cryptography (EC P-256 keys)  •  passlib + bcrypt (password hashing)",
        "• Altme Wallet v2.x+ (mobile credential wallet — iOS/Android)",
        "• Cloudflare Tunnel (cloudflared) for public HTTPS endpoint",
        "• W3C DID Core v1.0, W3C VC Data Model v1.1, OpenID4VCI Draft 13, OpenID4VP Draft 20",
    ],
    fill_color=RGBColor(0xE8, 0xEA, 0xF6), border_color=MED_BLUE,
    title_color=MED_BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 26 — Section Header: Acceptance & Traceability
# ══════════════════════════════════════════════════════════════════════════════

section_header("11. Acceptance Criteria & Traceability", "Verification Requirements and Requirement Mapping")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 27 — Acceptance Criteria
# ══════════════════════════════════════════════════════════════════════════════

s = content_slide("Acceptance Criteria")

ac_data = [
    ["Criterion", "Description", "Verification Method"],
    ["AC-01", "Student can register and receive JWT token", "E2E test: POST /auth/register → /auth/login"],
    ["AC-02", "Admin can issue credential offer with QR code", "E2E test: POST /admin/issue-vc → QR generated"],
    ["AC-03", "Wallet scans QR, exchanges code, receives signed VC", "Manual test with Altme wallet"],
    ["AC-04", "VC JWT has valid ES256 signature from issuer", "Unit test: sign + verify round-trip"],
    ["AC-05", "Verifier initiates OID4VP, wallet responds with VP", "Manual test with Altme wallet"],
    ["AC-06", "Backend verifies VP (holder sig + VC issuer sig)", "E2E test: POST callback → VERIFIED"],
    ["AC-07", "DID resolution works for did:web, did:jwk, did:key", "Unit test: resolve each method"],
    ["AC-08", "Key rotation archives old key, updates DID Document", "E2E test: rotate-key endpoint"],
    ["AC-09", "UI displays correct role-based views", "Manual test: login as each role"],
]
add_table(s, Inches(0.3), Inches(1.0), Inches(9.4), ac_data,
          col_widths_pct=[0.1, 0.52, 0.38])


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 28 — Requirements Traceability Matrix
# ══════════════════════════════════════════════════════════════════════════════

s = content_slide("Requirements Traceability Matrix")

trace_data = [
    ["Req ID", "Requirement", "Module / File", "Test / Acceptance"],
    ["FR-01", "Student registration", "auth_router.py", "AC-01"],
    ["FR-02", "JWT login", "auth_router.py, auth.py", "AC-01"],
    ["FR-05", "Create DID", "student_router.py", "AC-07"],
    ["FR-06", "Key rotation", "student_router.py", "AC-08"],
    ["FR-10", "Credential offer", "admin_router.py", "AC-02"],
    ["FR-11", "VC signing (OID4VCI)", "admin_router.py, crypto.py", "AC-03, AC-04"],
    ["FR-13", "OID4VP init", "verify_router.py", "AC-05"],
    ["FR-15", "VP callback verify", "verify_router.py, resolver.py", "AC-06"],
    ["FR-16", "Poll status", "verify_router.py", "AC-06"],
]
add_table(s, Inches(0.3), Inches(1.0), Inches(9.4), trace_data,
          col_widths_pct=[0.08, 0.28, 0.34, 0.3])

add_callout_box(s, Inches(0.3), Inches(3.8), Inches(9.4), Inches(1.2),
    "Traceability Note",
    [
        "All functional requirements map to specific router files and are covered by acceptance criteria.",
        "Non-functional requirements (Security, Performance) are addressed by architecture decisions (bcrypt, JWT, ES256, async FastAPI).",
    ])


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 29 — Section Header: Technology Stack
# ══════════════════════════════════════════════════════════════════════════════

section_header("12. Technology Stack", "Languages, Frameworks, Libraries, and Tools")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 30 — Technology Stack Content
# ══════════════════════════════════════════════════════════════════════════════

s = content_slide("Technology Stack")

tech_data = [
    ["Layer", "Technology", "Purpose"],
    ["Backend Framework", "FastAPI (Python 3.10+)", "Async web framework with auto-docs"],
    ["ORM", "SQLAlchemy 2.x", "Database abstraction layer"],
    ["Database", "SQLite (dev) / PostgreSQL (prod)", "Persistent data storage"],
    ["Authentication", "python-jose (JWT) + passlib (bcrypt)", "Token auth + password hashing"],
    ["Cryptography", "cryptography (P-256) + ES256", "Key generation, VC/VP signing"],
    ["Frontend", "Vanilla HTML/JS/CSS (SPA)", "Single-page app, no framework"],
    ["Wallet", "Altme (iOS/Android)", "OID4VCI + OID4VP compatible"],
    ["Tunnel", "Cloudflare Tunnel (cloudflared)", "Public HTTPS for wallet access"],
    ["Standards", "W3C DID, W3C VC, OID4VCI, OID4VP", "Interoperability protocols"],
]
add_table(s, Inches(0.3), Inches(1.0), Inches(9.4), tech_data,
          col_widths_pct=[0.2, 0.38, 0.42])

add_callout_box(s, Inches(0.3), Inches(4.0), Inches(9.4), Inches(1.1),
    "Project Structure",
    [
        "app/  →  main.py, auth.py, crypto.py, database.py, models.py, schemas.py, resolver.py",
        "app/routers/  →  auth_router.py, student_router.py, admin_router.py, verify_router.py, did_router.py",
        "static/  →  index.html, app.js, styles.css   |   test/  →  test_new_flow.py, test_e2e.py",
    ])


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 31 — Section Header: UI Design
# ══════════════════════════════════════════════════════════════════════════════

section_header("13. UI Design", "Role-Based Dashboards with Focus-Mode Navigation")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 32 — UI Screenshots / Wireframes — Login
# ══════════════════════════════════════════════════════════════════════════════

s = content_slide("UI Design — Login / Registration")

# Wireframe-style mockup
def wireframe_box(slide, left, top, w, h, label, sub_elements=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GRAY
    shape.line.color.rgb = RGBColor(0xBD, 0xBD, 0xBD)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.font.name = "Calibri"
    if sub_elements:
        for elem in sub_elements:
            p = tf.add_paragraph()
            p.text = elem
            p.font.size = Pt(7)
            p.font.color.rgb = RGBColor(0x75, 0x75, 0x75)
            p.font.name = "Calibri"

# Login wireframe
wireframe_box(s, Inches(0.5), Inches(1.1), Inches(4.0), Inches(3.8),
    "LOGIN SCREEN",
    [
        "┌─────────────────────────┐",
        "│  IIT Hyderabad Portal   │",
        "│  ═══════════════════    │",
        "│  [ Roll Number      ]   │",
        "│  [ ●●●●●●●●        ]   │",
        "│  [ Role: ▼ student  ]   │",
        "│                         │",
        "│  [ ▶ LOGIN ]            │",
        "│                         │",
        "│  Don't have an account? │",
        "│  [ Register → ]         │",
        "└─────────────────────────┘",
    ])

# Registration wireframe
wireframe_box(s, Inches(5.5), Inches(1.1), Inches(4.0), Inches(3.8),
    "REGISTRATION SCREEN",
    [
        "┌─────────────────────────┐",
        "│  Create Account         │",
        "│  ═══════════════════    │",
        "│  [ Full Name        ]   │",
        "│  [ Roll Number      ]   │",
        "│  [ ●●●●●●●●        ]   │",
        "│  [ Role: ▼ student  ]   │",
        "│                         │",
        "│  [ ▶ REGISTER ]         │",
        "│                         │",
        "│  Already registered?    │",
        "│  [ ← Login ]            │",
        "└─────────────────────────┘",
    ])


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 33 — UI — Admin Dashboard
# ══════════════════════════════════════════════════════════════════════════════

s = content_slide("UI Design — Admin (Issuer) Dashboard")

wireframe_box(s, Inches(0.3), Inches(1.1), Inches(9.4), Inches(4.0),
    "ADMIN DASHBOARD — Focus-Mode Shell",
    [
        "┌──┬─────────────────────────────────────────────────────────────────────┐",
        "│☰ │   Admin Dashboard  ▸  Students                    [ ⏻ Logout ]     │",
        "├──┤─────────────────────────────────────────────────────────────────────┤",
        "│  │                                                                     │",
        "│ MENU │   Student List — Overview & Credential Issuance                 │",
        "│  │   ┌────────┬───────────┬─────────┬───────┬────────────┐            │",
        "│ ☐ Overview │ Name  │ Roll No   │ DID?  │ Actions  │            │",
        "│ ☐ Students │       │           │       │          │            │",
        "│ ☐ Issue VC │ Sankalp│ 2021CS123 │  ✓    │ [Issue▼]│            │",
        "│  │         │ Rahul  │ 2021EE045 │  ✗    │ [—]     │            │",
        "│  │         └────────┴───────────┴───────┴────────────┘            │",
        "│  │                                                                     │",
        "│  │   Issue Credential: [ Select Student ▼ ] [ Type: Degree ▼ ]         │",
        "│  │   [ Degree ] [ Branch ] [ CGPA ] [ Year ] → [ Generate QR ]         │",
        "└──┴─────────────────────────────────────────────────────────────────────┘",
    ])


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 34 — UI — Student Dashboard
# ══════════════════════════════════════════════════════════════════════════════

s = content_slide("UI Design — Student Dashboard")

wireframe_box(s, Inches(0.3), Inches(1.1), Inches(9.4), Inches(4.0),
    "STUDENT DASHBOARD — Focus-Mode Shell",
    [
        "┌──┬─────────────────────────────────────────────────────────────────────┐",
        "│☰ │   Student Dashboard  ▸  My Credentials            [ ⏻ Logout ]     │",
        "├──┤─────────────────────────────────────────────────────────────────────┤",
        "│  │                                                                     │",
        "│ MENU │   My Credentials                                                │",
        "│  │   ┌─────────────────────────────────────────────────┐              │",
        "│ ☐ Overview │  PENDING OFFERS (scan QR with Altme)      │              │",
        "│ ☐ Credentials │  ╔═══════════════════════════╗         │              │",
        "│ ☐ Challenges  │  ║ B.Tech Degree  [QR CODE]  ║         │              │",
        "│  │            │  ║ Internship Cert [QR CODE]  ║         │              │",
        "│  │            │  ╚═══════════════════════════╝         │              │",
        "│  │   ├─────────────────────────────────────────────────┤              │",
        "│  │   │  CLAIMED CREDENTIALS (in wallet)                │              │",
        "│  │   │  ✓ UniversityDegreeCredential — 2025-01-15     │              │",
        "└──┴───┴─────────────────────────────────────────────────────────────────┘",
    ])


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 35 — UI — Verifier Dashboard
# ══════════════════════════════════════════════════════════════════════════════

s = content_slide("UI Design — Verifier Dashboard")

wireframe_box(s, Inches(0.3), Inches(1.1), Inches(9.4), Inches(4.0),
    "VERIFIER DASHBOARD — Focus-Mode Shell",
    [
        "┌──┬─────────────────────────────────────────────────────────────────────┐",
        "│☰ │   Verifier Dashboard  ▸  Proof Request             [ ⏻ Logout ]    │",
        "├──┤─────────────────────────────────────────────────────────────────────┤",
        "│  │                                                                     │",
        "│ MENU │   Request Credential Proof                                      │",
        "│  │   ┌─────────────────────────────────────────────┐                  │",
        "│ ☐ Request │  Select Student: [ ▼ did:web:...  ]    │                  │",
        "│ ☐ History │                                         │                  │",
        "│ ☐ Overview│  [ ▶ Generate QR ]                      │                  │",
        "│  │        │                                         │                  │",
        "│  │        │  ╔═══════════════════╗    Status:       │                  │",
        "│  │        │  ║   [QR CODE]       ║    ⏳ PENDING... │                  │",
        "│  │        │  ║  openid4vp://     ║    → ✓ VERIFIED │                  │",
        "│  │        │  ╚═══════════════════╝                  │                  │",
        "└──┴────────┴─────────────────────────────────────────────────────────────┘",
    ])


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 36 — Section Header: Reference Standards
# ══════════════════════════════════════════════════════════════════════════════

section_header("14. Reference Standards", "W3C, IETF, OpenID Foundation, ISO")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 37 — Reference Standards
# ══════════════════════════════════════════════════════════════════════════════

s = content_slide("Reference Standards & Specifications")

add_callout_box(s, Inches(0.3), Inches(1.0), Inches(4.5), Inches(1.8),
    "① W3C Standards (Define Quality)",
    [
        "• DID Core v1.0 — Decentralized Identifiers",
        "• VC Data Model v1.1 — Verifiable Credentials",
        "• DID:web Method Specification",
        "• JsonWebKey2020 Verification Method",
    ],
    fill_color=RGBColor(0xE8, 0xEA, 0xF6), border_color=MED_BLUE,
    title_color=MED_BLUE)

add_callout_box(s, Inches(5.2), Inches(1.0), Inches(4.5), Inches(1.8),
    "② OpenID Foundation (Build Quality)",
    [
        "• OID4VCI Draft 13+ — Credential Issuance",
        "• OID4VP Draft 20 — Verifiable Presentations",
        "• DIF Presentation Exchange v2.0",
        "• Pre-Authorized Code Grant Type",
    ],
    fill_color=RGBColor(0xE8, 0xF5, 0xE9), border_color=ACCENT_GREEN,
    title_color=ACCENT_GREEN)

add_callout_box(s, Inches(0.3), Inches(3.1), Inches(4.5), Inches(1.8),
    "③ IETF Standards (Security)",
    [
        "• RFC 7518 — JSON Web Algorithms (ES256)",
        "• RFC 7517 — JSON Web Key (JWK)",
        "• RFC 7519 — JSON Web Token (JWT)",
        "• RFC 6749 — OAuth 2.0 Framework",
    ],
    fill_color=RGBColor(0xFD, 0xF1, 0xE0), border_color=ACCENT_ORANGE,
    title_color=ACCENT_ORANGE)

add_callout_box(s, Inches(5.2), Inches(3.1), Inches(4.5), Inches(1.8),
    "④ ISO Quality Standards (Reference)",
    [
        "• ISO/IEC/IEEE 29148 — Requirements Eng.",
        "• ISO/IEC 25010:2011 — Product Quality Model",
        "• ISO/IEC 25019:2023 — Quality in Use",
        "• OWASP Top 10 — Security Best Practices",
    ],
    fill_color=RGBColor(0xF3, 0xE5, 0xF5), border_color=ACCENT_PURPLE,
    title_color=ACCENT_PURPLE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 38 — Crypto Flow Diagram
# ══════════════════════════════════════════════════════════════════════════════

s = content_slide("Cryptographic Flow — Key Management & Signing")

# Flow steps
steps_crypto = [
    ("1. Key Generation", "generate_ec_key_pair()\n→ EC P-256 private/public", Inches(0.3), Inches(1.1)),
    ("2. DID Construction", "did:web:{domain}:students:{id}\n+ publicKeyJwk in DID Doc", Inches(3.5), Inches(1.1)),
    ("3. VC Signing", "Issuer signs VC payload\nwith issuer private key (ES256)", Inches(6.7), Inches(1.1)),
    ("4. Wallet Binding", "Wallet proof JWT header\n→ extract wallet pub key\n→ update student DID doc", Inches(0.3), Inches(2.8)),
    ("5. VP Construction", "Wallet wraps VC in VP\nsigns with holder key (ES256)", Inches(3.5), Inches(2.8)),
    ("6. Verification", "Resolve DIDs → get pub keys\nVerify VP sig (holder)\nVerify VC sig (issuer)", Inches(6.7), Inches(2.8)),
]

colors_crypto = [DARK_BLUE, RGBColor(0x15, 0x65, 0xC0), ACCENT_GREEN,
                 ACCENT_PURPLE, ACCENT_ORANGE, RGBColor(0xC6, 0x28, 0x28)]

for i, (title, desc, x, y) in enumerate(steps_crypto):
    draw_block(s, x, y, Inches(2.8), Inches(0.4), title, colors_crypto[i], WHITE, 10)
    add_textbox(s, x + Inches(0.05), y + Inches(0.42), Inches(2.7), Inches(0.9),
                desc, font_size=9, color=DARK_GRAY)

# Flow arrows
draw_arrow_label(s, Inches(3.1), Inches(1.3), Inches(0.5), Inches(0.25), "→", 14)
draw_arrow_label(s, Inches(6.3), Inches(1.3), Inches(0.5), Inches(0.25), "→", 14)
draw_arrow_label(s, Inches(3.1), Inches(3.0), Inches(0.5), Inches(0.25), "→", 14)
draw_arrow_label(s, Inches(6.3), Inches(3.0), Inches(0.5), Inches(0.25), "→", 14)
draw_arrow_label(s, Inches(1.5), Inches(2.4), Inches(0.5), Inches(0.25), "↓", 14)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 39 — API Endpoints Summary
# ══════════════════════════════════════════════════════════════════════════════

s = content_slide("Complete API Endpoint Reference")

api_data = [
    ["Method", "Endpoint", "Auth", "Purpose"],
    ["POST", "/auth/register", "None", "Register new user"],
    ["POST", "/auth/login", "None", "Login, get JWT"],
    ["POST", "/auth/token", "None", "OID4VCI token exchange"],
    ["POST", "/students/did", "Student", "Create DID"],
    ["POST", "/students/did/rotate-key", "Student", "Rotate DID key"],
    ["GET", "/students/vcs", "Student", "List VCs + offers"],
    ["GET", "/students/challenges", "Student", "List verification challenges"],
    ["GET", "/admin/students", "Admin", "List all students"],
    ["POST", "/admin/issue-vc", "Admin", "Create credential offer"],
    ["POST", "/admin/issue", "Bearer", "OID4VCI credential endpoint"],
    ["POST", "/verify/init", "Verifier", "Start OID4VP session"],
    ["GET", "/api/verify/request/{id}", "None", "Serve signed request JWT"],
    ["POST", "/api/verify/callback/{id}", "None", "Wallet VP callback"],
    ["GET", "/verify/status/{id}", "Verifier", "Poll verification result"],
    ["GET", "/verify/history", "Verifier", "List past verifications"],
    ["GET", "/.well-known/openid-credential-issuer", "None", "OID4VCI discovery"],
]
add_table(s, Inches(0.15), Inches(0.95), Inches(9.7), api_data,
          col_widths_pct=[0.08, 0.33, 0.1, 0.49])


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 40 — Thank You / End Slide
# ══════════════════════════════════════════════════════════════════════════════

s = title_slide(
    "Thank You",
    "IIT Hyderabad — Self-Sovereign Identity Portal\n"
    "Decentralized Credentials for the Future of Education"
)
add_textbox(s, Inches(0.8), Inches(4.0), Inches(8.4), Inches(0.8),
            "Questions?",
            font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────

output_path = "IIT_Hyderabad_SSI_Portal_Presentation.pptx"
prs.save(output_path)
print(f"\n✅  Presentation saved: {output_path}")
print(f"   Total slides: {len(prs.slides)}")
